import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(r"c:\Users\wilki\OneDrive\Documentos\Trabajo de Curso")
CHART_DIR = ROOT / "_work" / "charts"
IMG_REF_DIR = ROOT / "Casa_de_Assados_Sofia_15_Imagens_Referencia_PR_HR_v2"

def resolve_img(pr_name, chart_fallback):
    p1 = IMG_REF_DIR / pr_name
    if p1.exists():
        return p1
    p2 = CHART_DIR / chart_fallback
    if p2.exists():
        return p2
    p3 = ROOT / "_work" / chart_fallback
    if p3.exists():
        return p3
    return p2

IMG_FIG02 = resolve_img("Figura_02_Fachada_Sinalizacao_Embalagens.png", "brand_mockup_sofia.jpg")
IMG_FIG03 = resolve_img("Figura_03_Cardapio_Balcao.png", "cardapio_impresso_sofia.jpg")
IMG_FIG04 = resolve_img("Figura_04_CRM_Sofia_Mobile_WhatsApp.png", "cardapio_whatsapp_sofia.jpg")
IMG_FIG09 = resolve_img("Figura_09_Assadoras_Giratorias_GLP.png", "equip1_asadora_gas.jpg")
IMG_FIG10 = resolve_img("Figura_10_Churrasqueira_Costela_Bafo.png", "equip2_churrasqueira_carvao.jpg")
IMG_FIG11 = resolve_img("Figura_11_Coifa_Industrial_AISI304.png", "equip3_coifa_industrial.jpg")
IMG_FIG12 = resolve_img("Figura_12_Freezer_Horizontal_510L.png", "equip4_freezer_horizontal.jpg")
IMG_FIG13 = resolve_img("Figura_13_Refrigerador_Vertical_4_Portas.png", "equip5_refrigerador_inox.jpg")
IMG_FIG14 = resolve_img("Figura_14_Mesa_Inox_Balanca.png", "equip6_bancada_balanca.jpg")
IMG_FIG15_PT = CHART_DIR / "planta_baixa_sofia_pt_hd.png"
IMG_FIG15_ES = CHART_DIR / "planta_baixa_sofia_es_hd.png"
IMG_FIG16 = resolve_img("Figura_16_Combo_Classico_Sofia.png", "combo1_classico_sofia.jpg")
IMG_FIG17 = resolve_img("Figura_17_Combo_Costela_Suprema.png", "combo2_costela_sofia.jpg")
IMG_FIG18 = resolve_img("Figura_18_Combo_Dueto_Sofia.png", "combo3_dueto_sofia.jpg")
IMG_FIG19 = resolve_img("Figura_19_Kit_Churrasco_Familia.png", "combo4_familia_sofia.jpg")
IMG_FIG20 = resolve_img("Figura_20_Conceito_Final.png", "anexo_casa_assados_sofia.png")

COLOR_NAVY = RGBColor(0x1F, 0x38, 0x64)      # #1F3864
COLOR_RED = RGBColor(0xC0, 0x39, 0x2B)       # #C0392B
COLOR_GOLD = RGBColor(0xD4, 0xAC, 0x0D)      # #D4AC0D
COLOR_DARK = RGBColor(0x2C, 0x3E, 0x50)      # #2C3E50
COLOR_BG_CARD = RGBColor(0xF4, 0xF7, 0xFA)   # #F4F7FA
COLOR_BORDER = RGBColor(0xDC, 0xE4, 0xEC)    # #DCE4EC
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)     # #FFFFFF
COLOR_GRAY = RGBColor(0x7F, 0x8C, 0x8D)      # #7F8C8D
COLOR_LIGHT_BLUE = RGBColor(0xE2, 0xEC, 0xF6)# #E2ECF6
COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)     # #27AE60

def set_slide_background(slide, color=COLOR_WHITE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="CASA DE ASSADOS SOFIA | TCC"):
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_NAVY
    top_bar.line.color.rgb = COLOR_NAVY
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_RED
    accent.line.color.rgb = COLOR_RED

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.12), Inches(12.133), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_GOLD
    p_cat.font.name = "Arial"
    p_cat.space_after = Pt(2)
    
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    p_title.font.name = "Arial"

def add_footer(slide, current_page=None, total_pages=18, lang="pt"):
    footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35))
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = COLOR_BG_CARD
    footer_bar.line.color.rgb = COLOR_BORDER
    
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.18), Inches(12.133), Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    
    if lang == "pt":
        p.text = "Colégio Excelência • Curso Técnico em Administração e Informática • Autor: Wilkin Barban Rosabal • 2026"
    else:
        p.text = "Colégio Excelência • Carrera Técnica en Administración e Informática • Autor: Wilkin Barban Rosabal • 2026"
        
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_GRAY
    p.font.name = "Arial"
    
    if current_page:
        tb_page = slide.shapes.add_textbox(Inches(11.5), Inches(7.18), Inches(1.2), Inches(0.3))
        tf_p = tb_page.text_frame
        p_num = tf_p.paragraphs[0]
        p_num.text = f"{current_page} / {total_pages}"
        p_num.alignment = PP_ALIGN.RIGHT
        p_num.font.size = Pt(9)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_NAVY

def add_card(slide, left_in, top_in, width_in, height_in, bg_color=COLOR_BG_CARD, border_color=COLOR_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card

def add_bullet_list(text_frame, items, font_size=11, text_color=COLOR_DARK, space_after=6):
    for i, item in enumerate(items):
        p = text_frame.add_paragraph() if i > 0 or len(text_frame.paragraphs[0].text) > 0 else text_frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.name = "Arial"
        p.space_after = Pt(space_after)

def generate_presentation(lang="pt"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # SLIDE 1: CAPA
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, COLOR_NAVY)
    
    # Top banner
    card_banner = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
    card_banner.fill.solid()
    card_banner.fill.fore_color.rgb = COLOR_RED
    card_banner.line.fill.background()
    
    tb_title = s1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(4.8))
    tf1 = tb_title.text_frame
    tf1.word_wrap = True
    
    p0 = tf1.paragraphs[0]
    p0.text = "COLÉGIO EXCELÊNCIA" if lang == "pt" else "COLÉGIO EXCELÊNCIA"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_GOLD
    p0.font.name = "Arial"
    p0.space_after = Pt(4)
    
    p_sub0 = tf1.add_paragraph()
    p_sub0.text = "CURSO TÉCNICO EM ADMINISTRAÇÃO E INFORMÁTICA" if lang == "pt" else "CARRERA TÉCNICA EN ADMINISTRACIÓN E INFORMÁTICA"
    p_sub0.font.size = Pt(12)
    p_sub0.font.color.rgb = COLOR_WHITE
    p_sub0.font.name = "Arial"
    p_sub0.space_after = Pt(28)
    
    p1 = tf1.add_paragraph()
    p1.text = "CASA DE ASSADOS SOFIA"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.font.name = "Arial"
    p1.space_after = Pt(12)
    
    p2 = tf1.add_paragraph()
    p2.text = (
        "Plano de Negócio para Implantação de uma Microempresa de Assados com Gestão por CRM em Curitiba - PR"
        if lang == "pt" else
        "Plan de Negocios para la Implantación de una Microempresa de Asados con Gestión por CRM en Curitiba - PR"
    )
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_LIGHT_BLUE
    p2.font.name = "Arial"
    p2.space_after = Pt(32)
    
    p3 = tf1.add_paragraph()
    p3.text = "Autor: Wilkin Barban Rosabal  |  Curitiba - PR, 2026" if lang == "pt" else "Autor: Wilkin Barban Rosabal  |  Curitiba - PR, 2026"
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_GOLD
    p3.font.name = "Arial"
    
    # ----------------------------------------------------
    # SLIDE 2: O PROBLEMA & A OPORTUNIDADE
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, COLOR_WHITE)
    add_header(s2, "Problema de Mercado e Oportunidade Estratégica" if lang == "pt" else "Problema de Mercado y Oportunidad Estratégica")
    add_footer(s2, 2, 18, lang)
    
    # Card 1: Problema
    add_card(s2, 0.8, 1.5, 5.6, 5.3, COLOR_BG_CARD, COLOR_BORDER)
    tb = s2.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🚨 O Dilema dos Assadores Convencionais" if lang == "pt" else "🚨 El Dilema de los Asaderos Convencionales"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    p.space_after = Pt(14)
    
    prob_items = [
        "Filas Extensas e Espera Caótica: Clientes aguardam de 30 a 50 min em pé aos domingos." if lang == "pt" else "Filas Extensas y Espera Caótica: Clientes esperan de 30 a 50 min de pie los domingos.",
        "Falta de Previsibilidade de Demanda: Produção 'no escuro' gerando sobra de carne ou falta de produto às 12h30." if lang == "pt" else "Falta de Previsibilidad: Producción 'a ciegas' generando sobras o quiebre de stock a las 12:30.",
        "Desperdício Severo de Insumos: Mermas de 12% a 18% em carnes e guarnições não comercializadas." if lang == "pt" else "Desperdicio Severo: Pérdidas del 12% al 18% en carnes y guarniciones no vendidas.",
        "Zero Relacionamento com o Cliente: Vendas 100% anônimas e passivas, sem histórico nem recompra guiada." if lang == "pt" else "Cero Relación con el Cliente: Ventas 100% anónimas y pasivas, sin historial ni retención."
    ]
    add_bullet_list(tf, prob_items, font_size=12, text_color=COLOR_DARK, space_after=10)

    # Card 2: Oportunidade
    add_card(s2, 6.9, 1.5, 5.6, 5.3, COLOR_LIGHT_BLUE, COLOR_NAVY)
    tb = s2.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.0), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡 A Solução: Gastronomia Programada" if lang == "pt" else "💡 La Solución: Gastronomía Programada"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(14)
    
    sol_items = [
        "Tradição do Churrasco de Domingo: Hábito cultural consolidado nas famílias curitibanas." if lang == "pt" else "Tradición del Asado de Domingo: Hábito cultural arraigado en las familias de Curitiba.",
        "Pré-Venda Ativa com CRM Casa de Assados Sofia: Encomendas antecipadas na sexta-feira via WhatsApp próprio." if lang == "pt" else "Preventa Activa con CRM Casa de Assados Sofia: Encargos anticipados los viernes vía WhatsApp propio.",
        "Retirada em 90s sem Fila: Janelas de atendimento fracionadas de 15 em 15 minutos." if lang == "pt" else "Retiro en 90s sin Filas: Franjas horarias fraccionadas de 15 en 15 minutos.",
        "Desperdício Reduzido para < 3%: Compras na CEASA calibradas na medida exata da demanda real." if lang == "pt" else "Desperdicio Reducido a < 3%: Compras en CEASA calibradas a la medida exacta de la demanda."
    ]
    add_bullet_list(tf, sol_items, font_size=12, text_color=COLOR_DARK, space_after=10)

    # ----------------------------------------------------
    # SLIDE 3: O CONCEITO DO NEGÓCIO & PROPOSTA DE VALOR
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, COLOR_WHITE)
    add_header(s3, "Conceito do Negócio e Proposta de Valor" if lang == "pt" else "Concepto del Negocio y Propuesta de Valor")
    add_footer(s3, 3, 18, lang)
    
    add_card(s3, 0.8, 1.5, 6.8, 5.3, COLOR_BG_CARD, COLOR_BORDER)
    tb = s3.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(6.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "CASA DE ASSADOS SOFIA LTDA."
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(12)
    
    concept_items = [
        ("Slogan Oficial: ", "“O verdadeiro sabor do domingo na mesa da sua família.”" if lang == "pt" else "“El auténtico sabor del domingo en la mesa de su familia.”"),
        ("Localização Estratégica: ", "Rua Deputado Pinheiro Júnior, 1380, Umbará, Curitiba - PR." if lang == "pt" else "Rua Deputado Pinheiro Júnior, 1380, Umbará, Curitiba - PR."),
        ("Formato Operacional: ", "Dark Kitchen / Takeaway com retirada expressa e delivery próprio em raio de 5 km." if lang == "pt" else "Dark Kitchen / Takeaway con retiro exprés y delivery propio en radio de 5 km."),
        ("Dias de Funcionamento: ", "Sábados, Domingos e TODOS os Feriados Nacionais, Estaduais e Municipais." if lang == "pt" else "Sábados, Domingos y TODOS los Feriados Nacionales, Provinciales y Municipales."),
        ("Enquadramento Jurídico: ", "Sociedade Limitada Unipessoal (SLU) no Simples Nacional (Microempresa)." if lang == "pt" else "Sociedad Limitada Unipersonal (SLU) en el Simples Nacional (Microempresa)."),
        ("Diferencial-Chave: ", "4 combos familiares padronizados + Agendamento inteligente via CRM Casa de Assados Sofia." if lang == "pt" else "4 combos familiares estandarizados + Agendamiento inteligente vía CRM Casa de Assados Sofia.")
    ]
    for bold_lead, text in concept_items:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = bold_lead
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_DARK
        p.space_after = Pt(8)

    # Right 3D Mockup
    if IMG_FIG20.exists():
        s3.shapes.add_picture(str(IMG_FIG20), Inches(8.0), Inches(1.5), Inches(4.5), Inches(5.3))

    # ----------------------------------------------------
    # SLIDE 4: O CRM CASA DE ASSADOS SOFIA COMO COLUNA VERTEBRAL
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, COLOR_WHITE)
    add_header(s4, "CRM Casa de Assados Sofia: A Coluna Vertebral Tecnológica do Negócio" if lang == "pt" else "CRM Casa de Assados Sofia: La Columna Vertebral Tecnológica del Negocio")
    add_footer(s4, 4, 18, lang)
    
    # Col 1: Software Próprio & Stack Livre
    add_card(s4, 0.8, 1.5, 3.7, 5.3, COLOR_BG_CARD, COLOR_NAVY)
    tb = s4.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💻 Desenvolvimento Próprio" if lang == "pt" else "💻 Desarrollo Propio In-House"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(10)
    c1_items = [
        "Criado pelo autor integrando Administração + Informática (Colégio Excelência)." if lang == "pt" else "Creado por el autor integrando Administración + Informática (Colégio Excelência).",
        "Domínio Dinâmico Grátis:\nhttps://casadeasados.duckdns.org/" if lang == "pt" else "Dominio Dinámico Gratis:\nhttps://casadeasados.duckdns.org/",
        "Certificado SSL Let's Encrypt (HTTPS seguro e criptografado)." if lang == "pt" else "Certificado SSL Let's Encrypt (HTTPS seguro y cifrado).",
        "Stack 100% Livre: Linux Ubuntu, PostgreSQL, Python/FastAPI e HTML5." if lang == "pt" else "Stack 100% Libre: Linux Ubuntu, PostgreSQL, Python/FastAPI y HTML5."
    ]
    add_bullet_list(tf, c1_items, font_size=11, space_after=8)

    # Col 2: Custo Hiper-Reduzido
    add_card(s4, 4.8, 1.5, 3.7, 5.3, COLOR_LIGHT_BLUE, COLOR_RED)
    tb = s4.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💰 Custo Quase Zero" if lang == "pt" else "💰 Costo Casi Cero"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    p.space_after = Pt(10)
    c2_items = [
        "Infraestrutura Total: R$ 50,00/mês (VPS Linux) + R$ 2,00/mês (IA DeepSeek V4 Flash) = R$ 52,00/mês." if lang == "pt" else "Infraestructura Total: R$ 50,00/mes (VPS Linux) + R$ 2,00/mes (IA DeepSeek V4 Flash) = R$ 52,00/mes.",
        "IA Sofia Integrada: 1.000 mensagens/mês por apenas R$ 2,00 com Prompt Caching." if lang == "pt" else "IA Sofia Integrada: 1.000 mensajes/mes por solo R$ 2,00 con Prompt Caching.",
        "Economia vs SaaS Comum: Sistemas pagos cobram R$ 300 a R$ 800/mês + taxas." if lang == "pt" else "Ahorro vs SaaS Común: Sistemas pagos cobran R$ 300 a R$ 800/mes + tasas.",
        "Sem Taxas de Marketplace: Evita 20% a 27% de comissão abusiva do iFood." if lang == "pt" else "Sin Comisiones de Marketplace: Evita el 20% al 27% abusivo de iFood."
    ]
    add_bullet_list(tf, c2_items, font_size=11, space_after=8)

    # Col 3: Impacto Operacional & LGPD
    add_card(s4, 8.8, 1.5, 3.7, 5.3, COLOR_BG_CARD, COLOR_GREEN)
    tb = s4.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚡ Impacto & Fidelização" if lang == "pt" else "⚡ Impacto y Fidelización"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    p.space_after = Pt(10)
    c3_items = [
        "Janelas de 15 min: Limita a 6 pedidos/intervalo, eliminando gargalos." if lang == "pt" else "Franjas de 15 min: Limita a 6 pedidos/intervalo, eliminando cuellos de botella.",
        "Painel KDS em Cozinha: Comandas em tempo real para os assadores." if lang == "pt" else "Panel KDS en Cocina: Comandas en tiempo real para los asadores.",
        "Segmentação RFM: Disparos para VIPs e reativação automática." if lang == "pt" else "Segmentación RFM: Mensajes para VIPs y reactivación automática.",
        "Conformidade LGPD: Opt-in / Opt-out nativo e base criptografada." if lang == "pt" else "Conformidad LGPD: Opt-in / Opt-out nativo y base cifrada."
    ]
    add_bullet_list(tf, c3_items, font_size=11, space_after=8)

    # ----------------------------------------------------
    # SLIDE 5: ENGENHARIA DE CARDÁPIO (OS 4 COMBOS)
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, COLOR_WHITE)
    add_header(s5, "Engenharia de Cardápio: Os 4 Combos Padronizados" if lang == "pt" else "Ingeniería de Menú: Los 4 Combos Estandarizados")
    add_footer(s5, 5, 18, lang)
    
    combos_info = [
        ("Combo 1: O Clássico" if lang == "pt" else "Combo 1: El Clásico", "R$ 69,90", "CMV: R$ 26,50 (62,1%)", IMG_FIG16, "1 Frango recheado inteiro assado + maionese caseira 300g + farofa crocante 250g. (3-4 pessoas)" if lang == "pt" else "1 Pollo relleno entero asado + mayonesa casera 300g + farofa crocante 250g. (3-4 personas)"),
        ("Combo 2: Costela Suprema" if lang == "pt" else "Combo 2: Costilla Suprema", "R$ 119,90", "CMV: R$ 48,00 (60,0%)", IMG_FIG17, "1,0kg de Costela bovina premium ao bafo (6h) + mandioca na manteiga + vinagrete e farofa. (4 pessoas)" if lang == "pt" else "1,0kg Costilla vacuna al vapor (6h) + mandioca a la manteca + vinagreta y farofa. (4 personas)"),
        ("Combo 3: Dueto Sofia" if lang == "pt" else "Combo 3: Dueto Sofia", "R$ 94,90", "CMV: R$ 36,00 (62,1%)", IMG_FIG18, "Meio frango dourado + 500g costelinha suína marinada + batatas rústicas e farofa. (3-4 pessoas)" if lang == "pt" else "Medio pollo dorado + 500g costilla de cerdo marinada + patatas rústicas y farofa. (3-4 personas)"),
        ("Combo 4: Kit Família" if lang == "pt" else "Combo 4: Kit Familia", "R$ 169,90", "CMV: R$ 68,00 (60,0%)", IMG_FIG19, "1 Frango inteiro + 700g costela + 4 linguiças artesanais + 4 pães de alho + maionese e farofa grande. (5-6 pessoas)" if lang == "pt" else "1 Pollo entero + 700g costilla + 4 chorizos parrilleros + 4 panes de ajo + mayonesa y farofa grande. (5-6 personas)")
    ]
    
    for i, (title, price, cmv_str, img_p, desc) in enumerate(combos_info):
        left_pos = 0.8 + i * 2.95
        add_card(s5, left_pos, 1.5, 2.8, 5.3, COLOR_BG_CARD, COLOR_BORDER)
        if img_p.exists():
            s5.shapes.add_picture(str(img_p), Inches(left_pos + 0.15), Inches(1.65), Inches(2.5), Inches(1.8))
            
        tb = s5.shapes.add_textbox(Inches(left_pos + 0.15), Inches(3.55), Inches(2.5), Inches(3.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY
        p.space_after = Pt(2)
        
        p = tf.add_paragraph()
        p.text = price
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_RED
        p.space_after = Pt(2)
        
        p = tf.add_paragraph()
        p.text = cmv_str
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_GREEN
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_DARK

    # ----------------------------------------------------
    # SLIDE 6: IDENTIDADE VISUAL & COMUNICAÇÃO DE MARCA
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, COLOR_WHITE)
    add_header(s6, "Identidade Visual, Slogan e Fachada Comercial" if lang == "pt" else "Identidad Visual, Slogan y Fachada Comercial")
    add_footer(s6, 6, 18, lang)
    
    add_card(s6, 0.8, 1.5, 6.0, 5.3, COLOR_BG_CARD, COLOR_BORDER)
    tb = s6.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.4), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Posicionamento Visual de Excelência" if lang == "pt" else "Posicionamiento Visual de Excelencia"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(12)
    
    brand_bullets = [
        "Paleta de Cores Institucional: Vermelho Brasa (#C0392B), Dourado Assado (#D4AC0D), Azul Confiança (#1F3864) e Preto Carvão (#2C3E50)." if lang == "pt" else "Paleta Institucional: Rojo Brasa (#C0392B), Dorado Asado (#D4AC0D), Azul Confianza (#1F3864) y Negro Carbón (#2C3E50).",
        "Fachada Moderna em Madeira e Aço: Revestimento rústico sofisticado com painel preto fosco e letreiro luminoso em acrílico 3D." if lang == "pt" else "Fachada Moderna en Madera y Acero: Revestimiento rústico sofisticado con panel negro mate y letrero luminoso en acrílico 3D.",
        "Embalagens Térmicas Seladas: Caixas e sacolas kraft personalizadas com lacre inviolável '100% Quente'." if lang == "pt" else "Envases Térmicos Sellados: Cajas y bolsas kraft personalizadas con precinto inviolable '100% Caliente'.",
        "Sinalização de Calçada: Cavalete rústico em madeira nobre (1,0m x 0,6m) com menu do dia e QR Code direto para o WhatsApp/CRM." if lang == "pt" else "Señalización de Acera: Caballete rústico en madera (1,0m x 0,6m) con menú del día y QR Code directo al WhatsApp/CRM."
    ]
    add_bullet_list(tf, brand_bullets, font_size=11.5, space_after=8)
    
    if IMG_FIG02.exists():
        s6.shapes.add_picture(str(IMG_FIG02), Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.3))

    # ----------------------------------------------------
    # SLIDE 7: CARDÁPIOS ESTRATÉGICOS (IMPRESSO E WHATSAPP)
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, COLOR_WHITE)
    add_header(s7, "Design dos Cardápios: Versão Impressa de Balcão e Digital WhatsApp" if lang == "pt" else "Diseño de Menús: Versión Impresa de Mostrador y Digital WhatsApp")
    add_footer(s7, 7, 18, lang)
    
    # Left Card: Impresso
    add_card(s7, 0.8, 1.5, 5.6, 5.3, COLOR_BG_CARD, COLOR_BORDER)
    tb = s7.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(5.2), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Cardápio Impresso para Balcão" if lang == "pt" else "Menú Impreso para Mostrador"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    if IMG_FIG03.exists():
        s7.shapes.add_picture(str(IMG_FIG03), Inches(1.1), Inches(2.25), Inches(5.0), Inches(4.3))
        
    # Right Card: WhatsApp
    add_card(s7, 6.9, 1.5, 5.6, 5.3, COLOR_BG_CARD, COLOR_BORDER)
    tb = s7.shapes.add_textbox(Inches(7.1), Inches(1.65), Inches(5.2), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Cardápio Interativo Mobile / WhatsApp" if lang == "pt" else "Menú Interactivo Mobile / WhatsApp"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    if IMG_FIG04.exists():
        s7.shapes.add_picture(str(IMG_FIG04), Inches(7.2), Inches(2.25), Inches(5.0), Inches(4.3))

    # ----------------------------------------------------
    # SLIDE 8: ANÁLISE DE MERCADO & PÚBLICO-ALVO
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, COLOR_WHITE)
    add_header(s8, "Análise de Mercado e Microrregião de Curitiba (Umbará)" if lang == "pt" else "Análisis de Mercado y Microrregión de Curitiba (Umbará)")
    add_footer(s8, 8, 18, lang)
    
    stats = [
        ("1.830.795", "Habitantes em Curitiba\n(IBGE 2024/2025)" if lang == "pt" else "Habitantes en Curitiba\n(IBGE 2024/2025)", COLOR_NAVY),
        ("R$ 120,06 Bi", "PIB Municipal de Curitiba\n(1ª da Região Sul / 6ª do Brasil)" if lang == "pt" else "PIB Municipal de Curitiba\n(1ª Región Sur / 6ª de Brasil)", COLOR_RED),
        ("R$ 67.691,30", "PIB per capita oficial\n(Alto poder aquisitivo)" if lang == "pt" else "PIB per cápita oficial\n(Alto poder adquisitivo)", COLOR_GOLD)
    ]
    for i, (val, lbl, col) in enumerate(stats):
        left_pos = 0.8 + i * 4.0
        add_card(s8, left_pos, 1.45, 3.7, 1.4, COLOR_BG_CARD, COLOR_BORDER)
        tb = s8.shapes.add_textbox(Inches(left_pos + 0.1), Inches(1.55), Inches(3.5), Inches(1.2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = col
        p.alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = lbl
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_DARK
        p.alignment = PP_ALIGN.CENTER

    add_card(s8, 0.8, 3.05, 5.8, 3.8, COLOR_BG_CARD, COLOR_BORDER)
    tb = s8.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(5.4), Inches(3.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📍 Microrregião Bairro Novo / Umbará" if lang == "pt" else "📍 Microrregión Bairro Novo / Umbará"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(8)
    mkt_bullets = [
        "População Regional: Mais de 165.000 habitantes nos bairros Umbará, Sítio Cercado e Ganchinho." if lang == "pt" else "Población Regional: Más de 165.000 habitantes en Umbará, Sítio Cercado y Ganchinho.",
        "Perfil Familiar: Classes B2, C1 e C2, famílias com 3 a 6 pessoas residentes em casas próprias e condomínios." if lang == "pt" else "Perfil Familiar: Estratos B2, C1 y C2, familias de 3 a 6 miembros en casas y condominios.",
        "Logística Privilegiada: Acesso rápido pela Rua Nicola Pellanda e Estrada do Ganchinho." if lang == "pt" else "Logística Privilegiada: Acceso rápido por Rua Nicola Pellanda y Estrada do Ganchinho."
    ]
    add_bullet_list(tf, mkt_bullets, font_size=11, space_after=6)

    add_card(s8, 6.9, 3.05, 5.6, 3.8, COLOR_WHITE, COLOR_BORDER)
    mix_chart = CHART_DIR / ("mix_pt.png" if lang == "pt" else "mix_es.png")
    if mix_chart.exists():
        s8.shapes.add_picture(str(mix_chart), Inches(7.05), Inches(3.15), Inches(5.3), Inches(3.6))

    # ----------------------------------------------------
    # SLIDE 9: PLANTA BAIXA & FLUXO SANITÁRIO (RDC 216)
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9, COLOR_WHITE)
    add_header(s9, "Planta Baixa Técnica e Fluxo Sanitário Unidirecional (RDC 216)" if lang == "pt" else "Plano Arquitectónico y Flujo Sanitario Unidireccional (RDC 216)")
    add_footer(s9, 9, 18, lang)
    
    add_card(s9, 0.8, 1.5, 4.8, 5.3, COLOR_BG_CARD, COLOR_BORDER)
    tb = s9.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(4.4), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Estrutura Física de 60 m² (10m x 6m)" if lang == "pt" else "Estructura Física de 60 m² (10m x 6m)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(10)
    
    layout_bullets = [
        "Conformidade Anvisa RDC 216: Fluxo linear unidirecional sem cruzamento entre cru e pronto." if lang == "pt" else "Conformidad Anvisa RDC 216: Flujo lineal unidireccional sin cruce entre crudo y cocido.",
        "Área de Cocção Isolada: Exaustão mecânica profissional com coifa em aço inox e dutos." if lang == "pt" else "Área de Cocción Aislada: Extracción mecánica profesional con campana inox y ductos.",
        "Superfícies Assépticas: Bancadas em aço inox AISI 304 e cubas gastronômicas GN." if lang == "pt" else "Superficies Asépticas: Mesadas en acero inox AISI 304 y cubas gastronómicas GN.",
        "Setor de Expedição Rápida: Balcão térmico de entrega expressa em < 90 segundos." if lang == "pt" else "Sector de Despacho Rápido: Mostrador térmico de entrega exprés en < 90 segundos."
    ]
    add_bullet_list(tf, layout_bullets, font_size=11, space_after=8)
    
    planta_img = IMG_FIG15_PT if lang == "pt" else IMG_FIG15_ES
    if planta_img.exists():
        s9.shapes.add_picture(str(planta_img), Inches(5.9), Inches(1.5), Inches(6.6), Inches(5.3))

    # ----------------------------------------------------
    # SLIDE 10: PARQUE DE EQUIPAMENTOS ADQUIRIDOS
    # ----------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10, COLOR_WHITE)
    add_header(s10, "Catálogo Fotográfico de Maquinários e Ativos Adquiridos" if lang == "pt" else "Catálogo Fotográfico de Maquinaria y Activos Adquiridos")
    add_footer(s10, 10, 18, lang)
    
    equip_items = [
        ("Asadoras a Gás GLP" if lang == "pt" else "Asadoras a Gas GLP", "2x Queimadores infravermelho" if lang == "pt" else "2x Quemadores infrarrojos", IMG_FIG09),
        ("Churrasqueira Carvão" if lang == "pt" else "Parrilla a Carbón", "Grelha elevatória e bafo 6h" if lang == "pt" else "Grelha elevadora y vapor 6h", IMG_FIG10),
        ("Coifa Industrial Inox" if lang == "pt" else "Campana Industrial Inox", "Exaustão mecânica (VISA)" if lang == "pt" else "Extracción mecánica (VISA)", IMG_FIG11),
        ("Freezer Horizontal 510L" if lang == "pt" else "Congelador Horizontal 510L", "Dupla ação (-18°C)" if lang == "pt" else "Doble acción (-18°C)", IMG_FIG12),
        ("Refrigerador Inox 4P" if lang == "pt" else "Refrigerador Inox 4P", "Vertical comercial (+2°C)" if lang == "pt" else "Vertical comercial (+2°C)", IMG_FIG13),
        ("Bancada Inox + Balança" if lang == "pt" else "Mesada Inox + Balanza", "AISI 304 e Inmetro digital" if lang == "pt" else "AISI 304 e Inmetro digital", IMG_FIG14)
    ]
    
    for i, (title, sub, img_p) in enumerate(equip_items):
        row = i // 3
        col = i % 3
        left_pos = 0.8 + col * 4.0
        top_pos = 1.45 + row * 2.75
        add_card(s10, left_pos, top_pos, 3.7, 2.55, COLOR_BG_CARD, COLOR_BORDER)
        if img_p.exists():
            s10.shapes.add_picture(str(img_p), Inches(left_pos + 0.15), Inches(top_pos + 0.15), Inches(1.8), Inches(2.25))
        tb = s10.shapes.add_textbox(Inches(left_pos + 2.05), Inches(top_pos + 0.3), Inches(1.5), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY
        p.space_after = Pt(4)
        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_DARK

    # ----------------------------------------------------
    # SLIDE 11: ESTRATÉGIA JURÍDICO-TRABALHISTA (DIARISTAS)
    # ----------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11, COLOR_WHITE)
    add_header(s11, "Estratégia Jurídico-Trabalhista e Segurança Operacional" if lang == "pt" else "Estrategia Jurídico-Laboral y Seguridad Operativa")
    add_footer(s11, 11, 18, lang)
    
    add_card(s11, 0.8, 1.5, 5.6, 5.3, COLOR_LIGHT_BLUE, COLOR_NAVY)
    tb = s11.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚖️ Contrato Intermitente (CLT 452-A)" if lang == "pt" else "⚖️ Contrato Intermitente (CLT 452-A)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(12)
    
    rh_bullets = [
        "Blindagem contra Passivo Trabalhista: Elimina riscos de vínculo clandestino ou ações na Justiça do Trabalho (TRT 9ª Região)." if lang == "pt" else "Blindaje contra Pasivos Laborales: Erradica riesgos de reclamos o juicios en la Justicia del Trabajo (TRT 9ª Región).",
        "Convocação Formal via CRM/WhatsApp: Pré-aviso de 72 horas com confirmação formal do trabalhador." if lang == "pt" else "Convocatoria Formal vía CRM/WhatsApp: Preaviso de 72 horas con aceptación expresa del trabajador.",
        "Pagamento Discriminado por Diária: Salário-base + DSR proporcional + 13º proporcional + Férias + 1/3 constitucional." if lang == "pt" else "Pago Discriminado por Jornal: Salario base + DSR + Aguinaldo proporcional + Vacaciones proporcionales + 1/3.",
        "Exames Ocupacionais (ASO): Atestados médicos admissionais com exames coprológicos para manipuladores de alimentos." if lang == "pt" else "Exámenes Ocupacionales (ASO): Aptitud médica y exámenes coprológicos obligatorios para manipuladores."
    ]
    add_bullet_list(tf, rh_bullets, font_size=11.5, space_after=8)
    
    add_card(s11, 6.9, 1.5, 5.6, 5.3, COLOR_BG_CARD, COLOR_BORDER)
    tb = s11.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Quadro de Funções e Diárias" if lang == "pt" else "Estructura de Funciones y Jornales"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    p.space_after = Pt(12)
    
    roles = [
        ("Sócio-Administrador", "Wilkin Barban", "Gestão geral, compras e CRM" if lang == "pt" else "Gestión general, compras y CRM"),
        ("Churrasqueiro Chefe", "R$ 120 / diária", "Cocção em grelha e bafo" if lang == "pt" else "Cocción en parrilla y vapor"),
        ("Auxiliar de Cozinha", "R$ 120 / diária", "Pré-preparo e guarnições" if lang == "pt" else "Preparación y guarniciones"),
        ("Auxiliar de Montagem", "R$ 120 / diária", "Embalagem e despacho KDS" if lang == "pt" else "Empaque y despacho KDS"),
        ("Entregador (Motoboy)", "R$ 120 / diária", "Delivery rápido raio 5km" if lang == "pt" else "Delivery rápido radio 5km")
    ]
    for r_title, r_pay, r_desc in roles:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {r_title} ({r_pay}): "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = COLOR_DARK
        r2 = p.add_run()
        r2.text = r_desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = COLOR_GRAY
        p.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 12: INVESTIMENTO INICIAL & FINANCIAMENTO
    # ----------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12, COLOR_WHITE)
    add_header(s12, "Investimento Inicial e Estrutura de Capital (R$ 38.000)" if lang == "pt" else "Inversión Inicial y Estructura de Capital (R$ 38.000)")
    add_footer(s12, 12, 18, lang)
    
    add_card(s12, 0.8, 1.5, 3.7, 5.3, COLOR_BG_CARD, COLOR_NAVY)
    tb = s12.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏗️ Investimento Fixo" if lang == "pt" else "🏗️ Inversión Fija"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(4)
    p_val = tf.add_paragraph()
    p_val.text = "R$ 24.500,00"
    p_val.font.size = Pt(22)
    p_val.font.bold = True
    p_val.font.color.rgb = COLOR_RED
    p_val.space_after = Pt(10)
    fix_items = [
        "2x Assadoras giratórias a gás: R$ 4.800" if lang == "pt" else "2x Asadoras a gas GLP: R$ 4.800",
        "1x Churrasqueira a carvão bafo: R$ 2.200" if lang == "pt" else "1x Parrilla al vapor: R$ 2.200",
        "1x Sistema de coifa industrial inox: R$ 4.200" if lang == "pt" else "1x Campana industrial inox: R$ 4.200",
        "1x Freezer horizontal 510L: R$ 3.100" if lang == "pt" else "1x Congelador 510L: R$ 3.100",
        "1x Refrigerador inox 4 portas: R$ 3.400" if lang == "pt" else "1x Refrigerador inox 4P: R$ 3.400",
        "2x Mesas inox + balança + utensílios: R$ 6.800" if lang == "pt" else "2x Mesadas inox + balanza: R$ 6.800"
    ]
    add_bullet_list(tf, fix_items, font_size=10.5, space_after=5)

    add_card(s12, 4.8, 1.5, 3.7, 5.3, COLOR_BG_CARD, COLOR_GOLD)
    tb = s12.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔄 Giro e Pré-Operacional" if lang == "pt" else "🔄 Capital de Trabajo"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(4)
    p_val = tf.add_paragraph()
    p_val.text = "R$ 13.500,00"
    p_val.font.size = Pt(22)
    p_val.font.bold = True
    p_val.font.color.rgb = COLOR_RED
    p_val.space_after = Pt(10)
    giro_items = [
        "Caução de Aluguel (3m) + 1º Mês: R$ 4.000" if lang == "pt" else "Garantía de Alquiler + 1º Mes: R$ 4.000",
        "Estoque Inicial de Insumos: R$ 2.500" if lang == "pt" else "Inventario Inicial Insumos: R$ 2.500",
        "Embalagens Térmicas Seladas: R$ 1.400" if lang == "pt" else "Envases Térmicos (1.000 un): R$ 1.400",
        "Licenças e Legalização (PMC/VISA): R$ 800" if lang == "pt" else "Licencias y Habilitación: R$ 800",
        "Fachada em Madeira e 3D: R$ 1.200" if lang == "pt" else "Cartel Fachada y Letrero 3D: R$ 1.200",
        "Marketing e Reserva de Liquidez: R$ 3.600" if lang == "pt" else "Marketing y Fondo de Reserva: R$ 3.600"
    ]
    add_bullet_list(tf, giro_items, font_size=10.5, space_after=5)

    add_card(s12, 8.8, 1.5, 3.7, 5.3, COLOR_LIGHT_BLUE, COLOR_NAVY)
    tb = s12.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏦 Estrutura de Financiamento" if lang == "pt" else "🏦 Estructura de Capital"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(4)
    p_val = tf.add_paragraph()
    p_val.text = "R$ 38.000,00"
    p_val.font.size = Pt(22)
    p_val.font.bold = True
    p_val.font.color.rgb = COLOR_NAVY
    p_val.space_after = Pt(10)
    fin_items = [
        "Capital Próprio: R$ 18.000,00 (47,37%) integralizado à vista pelo sócio." if lang == "pt" else "Capital Propio: R$ 18.000,00 (47,37%) integrado al contado por el socio.",
        "Microcrédito Fomento Paraná: R$ 20.000,00 (52,63%) via Banco do Empreendedor." if lang == "pt" else "Microcrédito Fomento Paraná: R$ 20.000,00 (52,63%) tasa subsidiada.",
        "Condições de Pagamento: 36 parcelas fixas de R$ 680,00/mês." if lang == "pt" else "Condiciones: 36 cuotas fijas de R$ 680,00/mes.",
        "Amortização: Totalmente inclusa nos custos fixos operacionais." if lang == "pt" else "Amortización: Totalmente contemplada en costos fijos."
    ]
    add_bullet_list(tf, fin_items, font_size=10.5, space_after=5)

    # ----------------------------------------------------
    # SLIDE 13: DEMONSTRATIVO DE RESULTADOS (DRE MENSAL)
    # ----------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_background(s13, COLOR_WHITE)
    add_header(s13, "Demonstrativo de Resultados do Exercício (DRE Base)" if lang == "pt" else "Estado de Resultados del Ejercicio (DRE Base)")
    add_footer(s13, 13, 18, lang)
    
    add_card(s13, 0.8, 1.5, 6.2, 5.3, COLOR_BG_CARD, COLOR_BORDER)
    tb = s13.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.8), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DRE Mensal - Cenário Base (160 Combos)" if lang == "pt" else "DRE Mensual - Escenario Base (160 Combos)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(10)
    
    dre_table = [
        ("(=) Receita Bruta Total (160 combos)", "R$ 15.809,00", "100,0%"),
        ("(-) Custos das Mercadorias (CMV)", "R$ 6.140,00", "38,8%"),
        ("(-) Impostos Simples Nacional (4,0%)", "R$ 632,36", "4,0%"),
        ("(-) Taxas de Meios de Pagamento (2,0%)", "R$ 316,18", "2,0%"),
        ("(=) Margem de Contribuição Bruta", "R$ 8.720,46", "55,2%"),
        ("(-) Custos Fixos Operacionais", "R$ 6.870,00", "43,5%"),
        ("(=) Lucro Operacional Líquido", "R$ 1.850,46", "11,7%")
    ] if lang == "pt" else [
        ("(=) Ingresos Brutos Totales (160 combos)", "R$ 15.809,00", "100,0%"),
        ("(-) Costo de Alimentos (CMV)", "R$ 6.140,00", "38,8%"),
        ("(-) Impuestos Simples Nacional (4,0%)", "R$ 632,36", "4,0%"),
        ("(-) Comisiones Medios de Pago (2,0%)", "R$ 316,18", "2,0%"),
        ("(=) Margen de Contribución Bruto", "R$ 8.720,46", "55,2%"),
        ("(-) Costos Fijos Operativos", "R$ 6.870,00", "43,5%"),
        ("(=) Utilidad Operativa Neta", "R$ 1.850,46", "11,7%")
    ]
    
    for label, val, pct in dre_table:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{label:42s}"
        r1.font.size = Pt(10.5)
        if "(=)" in label or "Lucro" in label or "Utilidad" in label:
            r1.font.bold = True
            r1.font.color.rgb = COLOR_NAVY if "Lucro" not in label and "Utilidad" not in label else COLOR_GREEN
        else:
            r1.font.color.rgb = COLOR_DARK
            
        r2 = p.add_run()
        r2.text = f"{val:>14s}  ({pct:>5s})"
        r2.font.size = Pt(10.5)
        r2.font.bold = True
        r2.font.color.rgb = COLOR_RED if "(-)" in label else (COLOR_GREEN if "Lucro" in label or "Utilidad" in label else COLOR_NAVY)
        p.space_after = Pt(4)

    add_card(s13, 7.3, 1.5, 5.2, 5.3, COLOR_WHITE, COLOR_BORDER)
    dre_chart = CHART_DIR / ("dre_pt.png" if lang == "pt" else "dre_es.png")
    if dre_chart.exists():
        s13.shapes.add_picture(str(dre_chart), Inches(7.45), Inches(1.65), Inches(4.9), Inches(5.0))

    # ----------------------------------------------------
    # SLIDE 14: PONTO DE EQUILÍBRIO & PAYBACK
    # ----------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_background(s14, COLOR_WHITE)
    add_header(s14, "Ponto de Equilíbrio Operacional e Payback" if lang == "pt" else "Punto de Equilibrio Operativo y Retorno de Inversión")
    add_footer(s14, 14, 18, lang)
    
    add_card(s14, 0.8, 1.5, 5.6, 5.3, COLOR_WHITE, COLOR_BORDER)
    be_chart = CHART_DIR / ("breakeven_pt.png" if lang == "pt" else "breakeven_es.png")
    if be_chart.exists():
        s14.shapes.add_picture(str(be_chart), Inches(0.95), Inches(1.65), Inches(5.3), Inches(5.0))
        
    add_card(s14, 6.7, 1.5, 5.8, 5.3, COLOR_BG_CARD, COLOR_BORDER)
    tb = s14.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(5.4), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Indicadores de Viabilidade e Margem" if lang == "pt" else "Indicadores de Viabilidad y Margen"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(12)
    
    kpis = [
        ("Índice Margem Contribuição:", " 55,16%", "Para cada R$ 100 vendidos, R$ 55,16 pagam custos e geram lucro." if lang == "pt" else "Por cada R$ 100 vendidos, R$ 55,16 pagan costos y dan ganancia."),
        ("Ponto de Equilíbrio (R$):", " R$ 12.454,37 / mês", "Faturamento mínimo mensal necessário para cobrir 100% dos custos." if lang == "pt" else "Facturación mínima mensual requerida para no tener pérdidas."),
        ("Ponto de Equilíbrio (Qtd):", " 126 combos / mês", "Venda de apenas 32 combos/fim de semana (16 por dia)." if lang == "pt" else "Venta de solo 32 combos/fin de semana (16 por día)."),
        ("Lucratividade Líquida:", " 11,71% no cenário base", "Margem saudável e protegida contra oscilações de insumos." if lang == "pt" else "Margen neto saludable y protegido frente a inflación."),
        ("Prazo de Retorno (Payback):", " 11 a 12 meses", "Recuperação integral dos R$ 38.000 investidos no primeiro ano." if lang == "pt" else "Recuperación completa de la inversión en el primer año.")
    ]
    for bold_t, val, sub in kpis:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = bold_t
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = val
        r2.font.bold = True
        r2.font.size = Pt(11)
        r2.font.color.rgb = COLOR_RED
        p.space_after = Pt(1)
        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = COLOR_GRAY
        p_sub.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 15: ALAVANCAGEM EM FERIADOS (2026-2028)
    # ----------------------------------------------------
    s15 = prs.slides.add_slide(blank_layout)
    set_slide_background(s15, COLOR_WHITE)
    add_header(s15, "Alavancagem Financeira em Feriados (Biênio 2026-2028)" if lang == "pt" else "Apalancamiento Financiero en Días Feriados (2026-2028)")
    add_footer(s15, 15, 18, lang)
    
    add_card(s15, 0.8, 1.5, 5.8, 5.3, COLOR_LIGHT_BLUE, COLOR_NAVY)
    tb = s15.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.4), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 Operação Estratégica em Feriados" if lang == "pt" else "🎯 Operación Estratégica en Feriados"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(10)
    
    fer_bullets = [
        "Conservadorismo Financeiro: A DRE base considera APENAS fins de semana (160 combos/mês)." if lang == "pt" else "Conservadurismo Contable: La DRE base considera SOLO fines de semana (160 combos/mes).",
        "Custos Fixos Já Amortizados: Aluguel, contador, internet e VPS já estão 100% pagos pelos fins de semana." if lang == "pt" else "Costos Fijos 100% Amortizados: Alquiler, contador, internet y VPS ya cubiertos por fines de semana.",
        "Margem Líquida Excedente: Quase 60% da receita de cada feriado vira lucro líquido puro." if lang == "pt" else "Margen Neto Excedente: Casi el 60% del ingreso de feriados se transforma en ganancia líquida.",
        "Reserva de Liquidez e Expansão: Os lucros extras blindam o caixa e aceleram quitação de empréstimos." if lang == "pt" else "Reserva y Expansión: Los beneficios extras blindan caja y permiten amortización anticipada."
    ]
    add_bullet_list(tf, fer_bullets, font_size=11.5, space_after=8)
    
    add_card(s15, 6.9, 1.5, 5.6, 5.3, COLOR_BG_CARD, COLOR_BORDER)
    tb = s15.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Projeção Consolidada dos Feriados" if lang == "pt" else "Proyección Consolidada de Feriados"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    p.space_after = Pt(12)
    
    f_stats = [
        ("Feriados Úteis Mapeados:", " 28 dias (2026-2028)", "Feriados de 2ª a 6ª feira em Curitiba." if lang == "pt" else "Feriados de lunes a viernes en Curitiba."),
        ("Volume de Vendas Extra:", " 625 combos", "Média de 20 a 25 combos por feriado." if lang == "pt" else "Promedio de 20 a 25 combos por feriado."),
        ("Receita Bruta Acumulada:", " + R$ 61.759,98", "Injeção extraordinária de faturamento." if lang == "pt" else "Inyección extraordinaria de ventas."),
        ("Lucro Líquido Adicional:", " + R$ 19.934,90", "Ganho líquido total transferido para a reserva de liquidez." if lang == "pt" else "Ganancia líquida total para reserva y expansión.")
    ]
    for lbl, val, sub in f_stats:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = lbl
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = COLOR_NAVY
        r2 = p.add_run()
        r2.text = val
        r2.font.bold = True
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = COLOR_GREEN
        p.space_after = Pt(1)
        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = COLOR_GRAY
        p_sub.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 16: ANÁLISE DE SENSIBILIDADE & MATRIZ SWOT
    # ----------------------------------------------------
    s16 = prs.slides.add_slide(blank_layout)
    set_slide_background(s16, COLOR_WHITE)
    add_header(s16, "Análise de Sensibilidade e Matriz Estratégica SWOT" if lang == "pt" else "Análisis de Sensibilidad y Matriz Estratégica FODA")
    add_footer(s16, 16, 18, lang)
    
    add_card(s16, 0.8, 1.5, 5.6, 5.3, COLOR_WHITE, COLOR_BORDER)
    scen_chart = CHART_DIR / ("scenarios_pt.png" if lang == "pt" else "scenarios_es.png")
    if scen_chart.exists():
        s16.shapes.add_picture(str(scen_chart), Inches(0.95), Inches(1.65), Inches(5.3), Inches(5.0))
        
    add_card(s16, 6.7, 1.5, 5.8, 5.3, COLOR_BG_CARD, COLOR_BORDER)
    tb = s16.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(5.4), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Matriz SWOT / FOFA Estratégica" if lang == "pt" else "Matriz FODA / SWOT Estratégica"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.space_after = Pt(10)
    
    swot_data = [
        ("FORÇAS (S): ", "Equipamentos novos, exaustão profissional, CRM próprio e cardápio enxuto." if lang == "pt" else "Equipos nuevos, extracción profesional, CRM propio y menú reducido."),
        ("OPORTUNIDADES (O): ", "Forte hábito de assados no Umbará, insatisfação com filas e abertura em feriados." if lang == "pt" else "Hábito arraigado de asados en Umbará, colas en competidores y feriados."),
        ("FRAQUEZAS (W): ", "Marca nova em fase de tração e capacidade física inicial limitada a 260 combos." if lang == "pt" else "Marca nueva en tracción y capacidad inicial acotada a 260 combos."),
        ("AMEAÇAS (T): ", "Oscilações de preço de carnes no atacado e clima chuvoso em fins de semana." if lang == "pt" else "Fluctuación de precios de carnes en mayoristas y mal clima.")
    ]
    for st, text in swot_data:
        p = tf.add_paragraph()
        r1 = p.add_run()
        r1.text = st
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = COLOR_RED if "FRAQUEZAS" in st or "AMEAÇAS" in st or "DEBILIDADES" in st or "AMENAZAS" in st else COLOR_NAVY
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = COLOR_DARK
        p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 17: GOVERNANÇA DE DADOS & SEGURANÇA (LGPD)
    # ----------------------------------------------------
    s17 = prs.slides.add_slide(blank_layout)
    set_slide_background(s17, COLOR_WHITE)
    add_header(s17, "Governança de Dados, Segurança e Conformidade com a LGPD" if lang == "pt" else "Gobernanza de Datos, Seguridad y Conformidad con la LGPD")
    add_footer(s17, 17, 18, lang)
    
    lgpd_cards = [
        ("🔐 Privacidade & LGPD" if lang == "pt" else "🔐 Privacidad y LGPD", [
            "Consentimento Explícito (Opt-in): Cliente autoriza receber mensagens via WhatsApp." if lang == "pt" else "Consentimiento Explícito (Opt-in): Cliente autoriza mensajes por WhatsApp.",
            "Opção de Saída (Opt-out): Descadastro imediato a qualquer momento digitando 'SAIR'." if lang == "pt" else "Opción de Baja (Opt-out): Cancelación inmediata enviando 'BAJA'.",
            "Minimização de Dados: Coleta apenas de Nome, Telefone e Endereço de entrega." if lang == "pt" else "Minimización de Datos: Solo Nombre, Teléfono y Dirección.",
            "Base Criptografada: Armazenamento seguro sem compartilhamento com terceiros." if lang == "pt" else "Base Cifrada: Almacenamiento seguro sin venta a terceros."
        ], COLOR_NAVY),
        ("🛡️ Segurança da Informação" if lang == "pt" else "🛡️ Seguridad de la Información", [
            "Tráfego 100% Criptografado: Certificado SSL/TLS Let's Encrypt (HTTPS)." if lang == "pt" else "Tráfico 100% Cifrado: Certificado SSL/TLS Let's Encrypt (HTTPS).",
            "Backups Diários Automatizados: Dump do PostgreSQL criptografado em AES-256." if lang == "pt" else "Copias de Respaldo Diarias: Dump PostgreSQL cifrado en AES-256.",
            "Recuperação em Minutos: RTO < 15 min e RPO < 24h em caso de falha de hardware." if lang == "pt" else "Recuperación Rápida: RTO < 15 min y RPO < 24h ante contingencias.",
            "Firewall e Acesso Restrito: Apenas portas essenciais (80/443) abertas no VPS." if lang == "pt" else "Firewall y Acceso Restringido: Solo puertos 80/443 abiertos en VPS."
        ], COLOR_GREEN)
    ]
    
    for i, (title, items, col) in enumerate(lgpd_cards):
        left_pos = 0.8 + i * 6.0
        add_card(s17, left_pos, 1.5, 5.7, 5.3, COLOR_BG_CARD, col)
        tb = s17.shapes.add_textbox(Inches(left_pos + 0.3), Inches(1.7), Inches(5.1), Inches(4.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(14)
        add_bullet_list(tf, items, font_size=12, text_color=COLOR_DARK, space_after=10)

    # ----------------------------------------------------
    # SLIDE 18: CONCLUSÃO & ENCERRAMENTO
    # ----------------------------------------------------
    s18 = prs.slides.add_slide(blank_layout)
    set_slide_background(s18, COLOR_NAVY)
    
    card_banner = s18.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
    card_banner.fill.solid()
    card_banner.fill.fore_color.rgb = COLOR_RED
    card_banner.line.fill.background()
    
    tb = s18.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "CONCLUSÃO & VIABILIDADE DO EMPREENDIMENTO" if lang == "pt" else "CONCLUSIÓN Y VIABILIDAD DEL EMPRENDIMIENTO"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(18)
    
    conc_items = [
        "Sinergia Interdisciplinar: Integração perfeita dos conhecimentos de Administração e Informática adquiridos no Colégio Excelência." if lang == "pt" else "Sinergia Interdisciplinaria: Integración de los conocimientos de Administración e Informática del Colégio Excelência.",
        "Viabilidade Econômica Plena: Ponto de equilíbrio de R$ 12.454,37 (~126 combos) abaixo do volume base (160 combos), com payback de 11 a 12 meses." if lang == "pt" else "Viabilidad Económica Plena: Punto de equilibrio de R$ 12.454,37 (~126 combos) inferior a la meta base (160 combos), con retorno en 11-12 meses.",
        "Diferencial Tecnológico Real: O CRM Casa de Assados Sofia resolve as dores clássicas de filas e desperdício de insumos com custo marginal de apenas R$ 50/mês." if lang == "pt" else "Diferencial Tecnológico Real: El CRM Casa de Assados Sofia resuelve las filas y el desperdicio de insumos con costo marginal de R$ 50/mes.",
        "Segurança Jurídica e Sanitária: Estrutura 100% blindada com CLT 452-A, RDC 216 Anvisa e alavancagem extraordinária em 28 feriados." if lang == "pt" else "Seguridad Jurídica y Sanitaria: Estructura blindada con CLT 452-A, RDC 216 Anvisa y apalancamiento extraordinario en 28 feriados."
    ]
    for it in conc_items:
        p = tf.add_paragraph()
        p.text = f"✔  {it}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_WHITE
        p.space_after = Pt(12)
        
    p_end = tf.add_paragraph()
    p_end.text = "\nMuito Obrigado! Aberto à Banca Examinadora para Perguntas." if lang == "pt" else "\n¡Muchas Gracias! Abierto al Tribunal Evaluador para Preguntas."
    p_end.font.size = Pt(15)
    p_end.font.bold = True
    p_end.font.color.rgb = COLOR_GOLD
    p_end.alignment = PP_ALIGN.CENTER
    
    out_file = ROOT / ("Apresentacao_Casa_de_Assados_Sofia_Portugues.pptx" if lang == "pt" else "Presentacion_Casa_de_Assados_Sofia_Espanol.pptx")
    prs.save(str(out_file))
    print(f"Presentation generated: {out_file} ({out_file.stat().st_size:,} bytes)")

if __name__ == "__main__":
    generate_presentation("pt")
    generate_presentation("es")
